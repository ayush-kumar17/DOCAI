"""
Document text extractors — one function per file type.

Each extractor returns a list of pages:
[
    {
        "page_number": 1,
        "text":        "extracted text here...",
        "section":     "Introduction",   # optional heading
        "has_table":   False,
        "has_image":   False,
    },
    ...
]

This standard format means the chunker doesn't care what
file type was uploaded — it always gets the same structure.
"""

import csv
import io
from pathlib import Path
from typing import List, Dict, Any

import pymupdf          # PyMuPDF — fast PDF text extraction
import pdfplumber       # PDF table extraction
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

from services.ocr.ocr_service import extract_text_from_image_bytes, extract_text_from_image_file
from utils.logger import get_logger

logger = get_logger(__name__)


# ──────────────────────────────────────────────
# PDF Extractor
# ──────────────────────────────────────────────

def extract_pdf(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from PDF page by page.

    Strategy:
    1. Use PyMuPDF to get text from each page (fast)
    2. If a page has < 50 chars, it's scanned → run OCR
    3. Use pdfplumber to detect and extract tables
    """
    pages = []
    path  = str(file_path)

    # Open with both libraries
    pdf_doc  = pymupdf.open(path)

    try:
        plumber_pdf = pdfplumber.open(path)
        plumber_pages = plumber_pdf.pages
    except Exception:
        plumber_pdf   = None
        plumber_pages = []

    for page_index, page in enumerate(pdf_doc):
        page_num = page_index + 1

        # ── Extract text ──────────────────────────
        text = page.get_text("text").strip()

        # ── OCR fallback for scanned pages ────────
        is_scanned = len(text) < 50
        if is_scanned:
            logger.info(f"Page {page_num} appears scanned — running OCR")
            try:
                # Render page to image at 200 DPI
                pix        = page.get_pixmap(dpi=200)
                img_bytes  = pix.tobytes("png")
                text       = extract_text_from_image_bytes(img_bytes)
            except Exception as e:
                logger.warning(f"OCR failed on page {page_num}: {e}")
                text = ""

        # ── Table extraction ──────────────────────
        tables_text = ""
        has_table   = False

        if plumber_pages and page_index < len(plumber_pages):
            try:
                plumber_page = plumber_pages[page_index]
                tables       = plumber_page.extract_tables()

                for table in tables:
                    if not table:
                        continue
                    has_table = True
                    rows = []
                    for row in table:
                        # Clean None values, join with pipe separator
                        clean_row = [str(cell).strip() if cell else "" for cell in row]
                        rows.append(" | ".join(clean_row))
                    tables_text += "\n[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]\n"

            except Exception as e:
                logger.warning(f"Table extraction failed on page {page_num}: {e}")

        # Combine prose text and table text
        full_text = (text + "\n" + tables_text).strip()

        if full_text:
            pages.append({
                "page_number": page_num,
                "text":        full_text,
                "section":     "",
                "has_table":   has_table,
                "has_image":   is_scanned,
            })

    pdf_doc.close()
    if plumber_pdf:
        plumber_pdf.close()

    logger.info(f"PDF extraction complete: {len(pages)} pages")
    return pages


# ──────────────────────────────────────────────
# DOCX Extractor
# ──────────────────────────────────────────────

def extract_docx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from Word documents.

    Strategy:
    - Walk paragraphs and detect headings (they become section names)
    - Group paragraphs between headings into one "page"
    - Extract tables separately
    """
    doc     = DocxDocument(file_path)
    pages   = []

    current_section = ""
    current_text    = []
    pseudo_page     = 1

    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ""
        text       = para.text.strip()

        if not text:
            continue

        # Headings become section markers and flush the current buffer
        if style_name.startswith("Heading"):
            if current_text:
                pages.append({
                    "page_number": pseudo_page,
                    "text":        "\n".join(current_text),
                    "section":     current_section,
                    "has_table":   False,
                    "has_image":   False,
                })
                pseudo_page += 1
                current_text = []

            current_section = text

        else:
            current_text.append(text)

    # Flush remaining text
    if current_text:
        pages.append({
            "page_number": pseudo_page,
            "text":        "\n".join(current_text),
            "section":     current_section,
            "has_table":   False,
            "has_image":   False,
        })
        pseudo_page += 1

    # Extract tables as separate pages
    for table in doc.tables:
        rows = []
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells
            )
            if row_text.strip():
                rows.append(row_text)

        if rows:
            table_text = "[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]"
            pages.append({
                "page_number": pseudo_page,
                "text":        table_text,
                "section":     current_section,
                "has_table":   True,
                "has_image":   False,
            })
            pseudo_page += 1

    logger.info(f"DOCX extraction complete: {len(pages)} sections")
    return pages


# ──────────────────────────────────────────────
# PPTX Extractor
# ──────────────────────────────────────────────

def extract_pptx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from PowerPoint, one slide = one page.
    """
    prs   = Presentation(file_path)
    pages = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts   = []
        title   = ""

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue

            # Detect title shape
            if shape.shape_type == 13 or (hasattr(shape, "placeholder_format")
                and shape.placeholder_format
                and shape.placeholder_format.idx == 0):
                title = text
            else:
                texts.append(text)

        full_text = "\n".join(texts)

        if title or full_text:
            pages.append({
                "page_number": slide_num,
                "text":        full_text,
                "section":     title or f"Slide {slide_num}",
                "has_table":   False,
                "has_image":   False,
            })

    logger.info(f"PPTX extraction complete: {len(pages)} slides")
    return pages


# ──────────────────────────────────────────────
# XLSX Extractor
# ──────────────────────────────────────────────

def extract_xlsx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract Excel spreadsheets.
    Each sheet becomes one page with the data as a table.
    Cap at 1000 rows per sheet to avoid huge chunks.
    """
    wb    = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    pages = []

    for sheet_num, sheet in enumerate(wb.worksheets, start=1):
        rows = []
        row_count = 0

        for row in sheet.iter_rows(values_only=True):
            if row_count >= 1000:
                rows.append("... (truncated at 1000 rows)")
                break

            clean = [str(cell) if cell is not None else "" for cell in row]
            if any(clean):   # skip completely empty rows
                rows.append(" | ".join(clean))
                row_count += 1

        if rows:
            table_text = f"[TABLE - Sheet: {sheet.title}]\n" + "\n".join(rows) + "\n[/TABLE]"
            pages.append({
                "page_number": sheet_num,
                "text":        table_text,
                "section":     f"Sheet: {sheet.title}",
                "has_table":   True,
                "has_image":   False,
            })

    wb.close()
    logger.info(f"XLSX extraction complete: {len(pages)} sheets")
    return pages


# ──────────────────────────────────────────────
# CSV Extractor
# ──────────────────────────────────────────────

def extract_csv(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract CSV files. Entire file becomes one table page.
    Cap at 2000 rows.
    """
    rows = []

    with open(file_path, "r", encoding="utf-8-sig", errors="ignore") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i >= 2000:
                rows.append("... (truncated at 2000 rows)")
                break
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))

    if not rows:
        return []

    table_text = "[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]"
    return [{
        "page_number": 1,
        "text":        table_text,
        "section":     "Data",
        "has_table":   True,
        "has_image":   False,
    }]


# ──────────────────────────────────────────────
# TXT / MD Extractor
# ──────────────────────────────────────────────

def extract_text_file(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract plain text or markdown files.
    Split into pseudo-pages of ~2000 characters each
    so very long files don't create one giant chunk.
    """
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    if not content.strip():
        return []

    # Split into ~2000 char chunks as pseudo-pages
    page_size = 2000
    pages     = []

    for i in range(0, len(content), page_size):
        chunk = content[i : i + page_size].strip()
        if chunk:
            pages.append({
                "page_number": (i // page_size) + 1,
                "text":        chunk,
                "section":     "",
                "has_table":   False,
                "has_image":   False,
            })

    logger.info(f"Text extraction complete: {len(pages)} pseudo-pages")
    return pages


# ──────────────────────────────────────────────
# Image Extractor
# ──────────────────────────────────────────────

def extract_image(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from image files using OCR.
    """
    logger.info(f"Running OCR on image: {file_path}")
    text = extract_text_from_image_file(file_path)

    if not text:
        return []

    return [{
        "page_number": 1,
        "text":        text,
        "section":     "",
        "has_table":   False,
        "has_image":   True,
    }]


# ──────────────────────────────────────────────
# Router — pick the right extractor
# ──────────────────────────────────────────────

def extract_document(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    """
    Main entry point — routes to the correct extractor by file type.
    """
    extractors = {
        "pdf":  extract_pdf,
        "docx": extract_docx,
        "pptx": extract_pptx,
        "xlsx": extract_xlsx,
        "csv":  extract_csv,
        "txt":  extract_text_file,
        "md":   extract_text_file,
        "png":  extract_image,
        "jpg":  extract_image,
        "jpeg": extract_image,
        "webp": extract_image,
        "tiff": extract_image,
    }

    extractor = extractors.get(file_type.lower())
    if not extractor:
        raise ValueError(f"No extractor for file type: {file_type}")

    return extractor(file_path)